import io
import PyPDF2
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text.strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines).strip()


def extract_text_from_txt(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            return file_bytes.decode(encoding).strip()
        except (UnicodeDecodeError, LookupError):
            continue
    return file_bytes.decode("utf-8", errors="replace").strip()


def extract_text(file_bytes: bytes, file_name: str) -> str:
    name_lower = file_name.lower()
    if name_lower.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif name_lower.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif name_lower.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif name_lower.endswith(".txt"):
        return extract_text_from_txt(file_bytes)
    else:
        raise ValueError(f"Qo'llab-quvvatlanmaydigan fayl turi: {file_name}")
